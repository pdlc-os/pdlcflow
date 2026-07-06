"""Entity CRUD — domains, squads, initiatives, repositories, projects.

Org-scoped (the org comes from the JWT when auth is on, else an `org_id` query
param). Writes go through `set_org_context`, so RLS lets a tenant only ever touch
its own rows. Repo tokens are stored via the secrets backend (never returned).
"""

from __future__ import annotations

import base64
import re
import uuid
from datetime import UTC, datetime

import httpx
from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile
from pdlc_graph.ports import (
    list_artifacts,
    put_artifact,
    read_artifact,
    reset_current_org,
    set_current_org,
)
from pydantic import BaseModel
from sqlalchemy import text

from ..auth.local import Identity, get_principal, resolve_org
from ..config import settings
from ..db.rls import set_org_context
from ..db.session import get_sync_engine
from ..secretstore import get_secrets

router = APIRouter(tags=["entities"])


def current_org(label: str):
    def _dep(
        org_id: str | None = Query(None),
        principal: Identity | None = Depends(get_principal),
    ) -> str:
        return resolve_org(principal, org_id, label)

    return _dep


def _slug(name: str) -> str:
    base = re.sub(r"[^a-z0-9]+", "-", name.lower()).strip("-") or "item"
    return f"{base[:40]}-{uuid.uuid4().hex[:4]}"


def _rows(org: str, sql: str, params: dict | None = None) -> list[dict]:
    with get_sync_engine(settings).begin() as c:
        set_org_context(c, org)
        return [dict(r) for r in c.execute(text(sql), params or {}).mappings().all()]


def _one(org: str, sql: str, params: dict, *, ensure_org: bool = False) -> dict:
    with get_sync_engine(settings).begin() as c:
        set_org_context(c, org)
        if ensure_org:
            # Self-host: make sure the tenant row exists (organizations has no RLS).
            c.execute(
                text("insert into organizations (id, name, slug, created_at, settings) "
                     "values (cast(:o as uuid), 'Default', 'org-' || left(cast(:o as text), 8), now(), '{}'::jsonb) "
                     "on conflict (id) do nothing"),
                {"o": org},
            )
        return dict(c.execute(text(sql), params).mappings().one())


# ---------------- Domains ----------------
class DomainCreate(BaseModel):
    name: str


@router.get("/domains")
def list_domains(org: str = Depends(current_org("/domains"))) -> dict:
    return {"domains": _rows(org, "select id, name from domains where org_id = :o order by name", {"o": org})}


@router.post("/domains")
def create_domain(body: DomainCreate, org: str = Depends(current_org("/domains"))) -> dict:
    return _one(org, "insert into domains (id, org_id, name) values (gen_random_uuid(), :o, :n) returning id, name",
                {"o": org, "n": body.name}, ensure_org=True)


# ---------------- Squads ----------------
class SquadCreate(BaseModel):
    name: str
    domain_id: str | None = None


@router.get("/squads")
def list_squads(org: str = Depends(current_org("/squads"))) -> dict:
    return {"squads": _rows(org, "select id, name, slug, domain_id from squads where org_id = :o order by name", {"o": org})}


@router.post("/squads")
def create_squad(body: SquadCreate, org: str = Depends(current_org("/squads"))) -> dict:
    return _one(org,
                "insert into squads (id, org_id, name, slug, domain_id) "
                "values (gen_random_uuid(), :o, :n, :s, :d) returning id, name, slug, domain_id",
                {"o": org, "n": body.name, "s": _slug(body.name), "d": body.domain_id}, ensure_org=True)


# ---------------- Initiatives ----------------
class InitiativeCreate(BaseModel):
    name: str
    status: str = "active"


@router.get("/initiatives")
def list_initiatives(org: str = Depends(current_org("/initiatives"))) -> dict:
    return {"initiatives": _rows(org, "select id, name, status from initiatives where org_id = :o order by name", {"o": org})}


@router.post("/initiatives")
def create_initiative(body: InitiativeCreate, org: str = Depends(current_org("/initiatives"))) -> dict:
    return _one(org,
                "insert into initiatives (id, org_id, name, status, created_at) "
                "values (gen_random_uuid(), :o, :n, :st, now()) returning id, name, status",
                {"o": org, "n": body.name, "st": body.status}, ensure_org=True)


# ---------------- Repositories ----------------
class RepoCreate(BaseModel):
    squad_id: str
    name: str
    url: str
    token: str | None = None
    default_branch: str = "main"
    provider: str = "github"


@router.get("/repositories")
def list_repositories(
    org: str = Depends(current_org("/repositories")),
    squad_id: str | None = Query(None),
) -> dict:
    sql = ("select id, name, url, squad_id, default_branch, provider, "
           "(token_secret_ref is not null) as has_token from repositories where org_id = :o")
    params: dict = {"o": org}
    if squad_id:
        sql += " and squad_id = :sq"
        params["sq"] = squad_id
    return {"repositories": _rows(org, sql + " order by name", params)}


@router.post("/repositories")
def create_repository(body: RepoCreate, org: str = Depends(current_org("/repositories"))) -> dict:
    rid = str(uuid.uuid4())
    ref = get_secrets().put(body.token, hint=rid) if body.token else None
    return _one(org,
                "insert into repositories (id, org_id, squad_id, name, url, default_branch, provider, token_secret_ref, created_at) "
                "values (:id, :o, :sq, :n, :u, :br, :p, :ref, now()) "
                "returning id, name, url, squad_id, default_branch, provider, (token_secret_ref is not null) as has_token",
                {"id": rid, "o": org, "sq": body.squad_id, "n": body.name, "u": body.url,
                 "br": body.default_branch, "p": body.provider, "ref": ref})


@router.delete("/repositories/{repo_id}")
def delete_repository(repo_id: str, org: str = Depends(current_org("/repositories"))) -> dict:
    rows = _rows(org, "delete from repositories where id = :id returning id", {"id": repo_id})
    if not rows:
        raise HTTPException(status_code=404, detail="repository not found")
    return {"deleted": repo_id}


# ---------------- Repo file browsing (GitHub) — powers the repo-backed memory panel ----------------
def _gh_owner_repo(url: str) -> tuple[str, str]:
    m = re.search(r"github\.com[:/]+([^/]+)/([^/]+?)(?:\.git)?/?$", url)
    if not m:
        raise HTTPException(status_code=400, detail="only github.com repositories are supported for file browsing")
    return m.group(1), m.group(2)


def _repo_row(org: str, repo_id: str) -> dict:
    rows = _rows(org,
                 "select url, provider, default_branch, token_secret_ref from repositories where id = :id",
                 {"id": repo_id})
    if not rows:
        raise HTTPException(status_code=404, detail="repository not found")
    return rows[0]


def _gh_get(org: str, repo_id: str, path: str):
    r = _repo_row(org, repo_id)
    if r["provider"] != "github":
        raise HTTPException(status_code=400, detail="file browsing currently supports github repos")
    owner, name = _gh_owner_repo(r["url"])
    token = get_secrets().resolve(r["token_secret_ref"])
    headers = {"Accept": "application/vnd.github+json"}
    if token:
        headers["Authorization"] = f"Bearer {token}"
    api = f"https://api.github.com/repos/{owner}/{name}/contents/{path.strip('/')}"
    try:
        resp = httpx.get(api, headers=headers, params={"ref": r["default_branch"]}, timeout=15)
    except Exception as exc:  # pragma: no cover - network
        raise HTTPException(status_code=502, detail=f"github unreachable: {exc}") from exc
    return resp


@router.get("/repositories/{repo_id}/files")
def repo_files(
    repo_id: str,
    org: str = Depends(current_org("/repositories")),
    path: str = Query(""),
) -> dict:
    resp = _gh_get(org, repo_id, path)
    if resp.status_code == 404:
        return {"path": path, "entries": []}
    if resp.status_code >= 400:
        raise HTTPException(status_code=502, detail=f"github error {resp.status_code}")
    data = resp.json()
    items = data if isinstance(data, list) else [data]
    return {"path": path, "entries": [
        {"name": x["name"], "path": x["path"], "type": x["type"], "size": x.get("size", 0)} for x in items]}


@router.get("/repositories/{repo_id}/file")
def repo_file(
    repo_id: str,
    org: str = Depends(current_org("/repositories")),
    path: str = Query(...),
) -> dict:
    resp = _gh_get(org, repo_id, path)
    if resp.status_code >= 400:
        raise HTTPException(status_code=502, detail=f"github error {resp.status_code}")
    data = resp.json()
    if isinstance(data, list):
        raise HTTPException(status_code=400, detail="path is a directory")
    content = data.get("content", "")
    if data.get("encoding") == "base64":
        content = base64.b64decode(content).decode("utf-8", "replace")
    return {"path": path, "name": data.get("name"), "content": content}


# ---------------- Projects (server-backed; supersedes the Studio client registry) ----------------
class ProjectCreate(BaseModel):
    name: str
    squad_id: str
    repository_id: str | None = None


@router.get("/projects")
def list_projects(org: str = Depends(current_org("/projects"))) -> dict:
    return {"projects": _rows(org,
            "select id, name, slug, squad_id, repository_id from projects where org_id = :o order by created_at desc",
            {"o": org})}


_MAX_UPLOAD = 15_000_000  # 15 MB


def _extract_doc_text(raw: bytes, name: str) -> str | None:
    """Pull text out of common binary docs so attachments reach the LLM."""
    import io
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    try:
        if ext == "pdf":
            from pypdf import PdfReader
            return "\n".join((p.extract_text() or "") for p in PdfReader(io.BytesIO(raw)).pages)
        if ext == "docx":
            import docx
            return "\n".join(p.text for p in docx.Document(io.BytesIO(raw)).paragraphs)
        if ext in ("xlsx", "xlsm"):
            from openpyxl import load_workbook
            out: list[str] = []
            for ws in load_workbook(io.BytesIO(raw), read_only=True, data_only=True).worksheets:
                out.append(f"# {ws.title}")
                out += ["\t".join("" if c is None else str(c) for c in row)
                        for row in ws.iter_rows(values_only=True)]
            return "\n".join(out)
        if ext == "pptx":
            from pptx import Presentation
            out2: list[str] = []
            for i, slide in enumerate(Presentation(io.BytesIO(raw)).slides, 1):
                out2.append(f"# Slide {i}")
                out2 += [s.text_frame.text for s in slide.shapes if s.has_text_frame]
            return "\n".join(out2)
    except Exception:  # pragma: no cover - malformed doc
        return None
    return None


@router.post("/uploads")
async def upload_file(
    file: UploadFile = File(...),
    project_id: str = Form(...),
    conversation_id: str = Form(...),
    org: str = Depends(current_org("/uploads")),
) -> dict:
    """Store a chat attachment under the tenant's artifact namespace, in a folder
    per project + conversation, with a timestamped filename (so re-uploading the
    same name in the same conversation never overwrites). Text-like files and
    extractable docs (pdf/docx/xlsx/pptx) return their text so the agent can use it.
    """
    raw = await file.read()
    if len(raw) > _MAX_UPLOAD:
        raise HTTPException(status_code=413, detail="file too large (max 15 MB)")
    name = (file.filename or "upload").replace("/", "_").replace("\\", "_")
    conv = re.sub(r"[^A-Za-z0-9_.-]", "", conversation_id) or "unfiled"
    stamp = datetime.now(UTC).strftime("%Y%m%dT%H%M%S")
    stored = f"{stamp}-{uuid.uuid4().hex[:6]}-{name}"  # timestamp + nonce => no overwrite

    try:
        text_content: str | None = raw.decode("utf-8")
        is_text = True
    except UnicodeDecodeError:
        text_content, is_text = None, False
    if text_content is None:
        text_content = _extract_doc_text(raw, name)  # pdf/docx/xlsx/pptx

    tok = set_current_org(org)
    try:
        if is_text:
            uri = put_artifact(project_id, f"uploads/{conv}/{stored}", raw.decode("utf-8"))
        else:
            uri = put_artifact(project_id, f"uploads/{conv}/{stored}.b64", base64.b64encode(raw).decode())
    finally:
        reset_current_org(tok)

    return {"id": uuid.uuid4().hex[:12], "filename": name, "stored_as": stored,
            "conversation_id": conv, "size": len(raw), "content_type": file.content_type,
            "is_text": is_text, "uri": uri, "text": (text_content[:20000] if text_content else None)}


@router.get("/projects/{project_id}/tasks")
def list_project_tasks(
    project_id: str, org: str = Depends(current_org("/projects"))
) -> dict:
    """The project's tasks (bd-NN) for the roadmap board. Org-scoped."""
    from pdlc_graph.ports import get_task_store

    rows = get_task_store().list(org, project_id)
    tasks = [
        {"external_id": r.get("external_id"), "title": r.get("title"),
         "status": r.get("status") or "open", "labels": r.get("labels") or [],
         "branch": r.get("branch"), "claimed_by": r.get("claimed_by")}
        for r in rows
    ]
    return {"project_id": project_id, "tasks": tasks}


@router.get("/projects/{project_id}/artifacts")
def list_project_artifacts(
    project_id: str, org: str = Depends(current_org("/projects"))
) -> dict:
    """Relative paths of the project's stored artifacts (PRD/design/decisions/
    deployments/uploads/…). Org-scoped via the turn context."""
    tok = set_current_org(org)
    try:
        paths = list_artifacts(project_id)
    finally:
        reset_current_org(tok)
    return {"project_id": project_id, "artifacts": paths}


@router.get("/projects/{project_id}/artifacts/content")
def get_project_artifact(
    project_id: str, path: str, org: str = Depends(current_org("/projects"))
) -> dict:
    """Fetch one artifact's text by relative path."""
    tok = set_current_org(org)
    try:
        try:
            content = read_artifact(project_id, path)
        except (KeyError, FileNotFoundError):
            raise HTTPException(status_code=404, detail="artifact not found") from None
    finally:
        reset_current_org(tok)
    return {"project_id": project_id, "path": path, "content": content}


@router.post("/projects")
def create_project(body: ProjectCreate, org: str = Depends(current_org("/projects"))) -> dict:
    return _one(org,
                "insert into projects (id, org_id, squad_id, repository_id, name, slug, created_at) "
                "values (gen_random_uuid(), :o, :sq, :r, :n, :s, now()) returning id, name, slug, squad_id, repository_id",
                {"o": org, "sq": body.squad_id, "r": body.repository_id, "n": body.name, "s": _slug(body.name)},
                ensure_org=True)


# ---------------- Rename + delete (org-scoped via RLS) ----------------
class RenameBody(BaseModel):
    name: str


def _rename(org: str, table: str, entity_id: str, name: str, returning: str) -> dict:
    rows = _rows(org, f"update {table} set name = :n where id = :id returning {returning}",
                 {"n": name, "id": entity_id})
    if not rows:
        raise HTTPException(status_code=404, detail=f"{table[:-1]} not found")
    return rows[0]


def _tx(org: str):
    """A single RLS-scoped transaction for multi-statement deletes."""
    eng = get_sync_engine(settings)
    ctx = eng.begin()
    conn = ctx.__enter__()
    set_org_context(conn, org)
    return ctx, conn


@router.patch("/domains/{entity_id}")
def rename_domain(entity_id: str, body: RenameBody, org: str = Depends(current_org("/domains"))) -> dict:
    return _rename(org, "domains", entity_id, body.name, "id, name")


@router.delete("/domains/{entity_id}")
def delete_domain(entity_id: str, org: str = Depends(current_org("/domains"))) -> dict:
    # squads.domain_id is ON DELETE SET NULL — squads survive, just un-grouped.
    if not _rows(org, "delete from domains where id = :id returning id", {"id": entity_id}):
        raise HTTPException(status_code=404, detail="domain not found")
    return {"deleted": entity_id}


@router.patch("/squads/{entity_id}")
def rename_squad(entity_id: str, body: RenameBody, org: str = Depends(current_org("/squads"))) -> dict:
    return _rename(org, "squads", entity_id, body.name, "id, name, slug, domain_id")


@router.delete("/squads/{entity_id}")
def delete_squad(entity_id: str, org: str = Depends(current_org("/squads"))) -> dict:
    # Cascades to the squad's projects (→ tasks/memory/gates), repositories, and links.
    if not _rows(org, "delete from squads where id = :id returning id", {"id": entity_id}):
        raise HTTPException(status_code=404, detail="squad not found")
    return {"deleted": entity_id}


@router.patch("/initiatives/{entity_id}")
def rename_initiative(entity_id: str, body: RenameBody, org: str = Depends(current_org("/initiatives"))) -> dict:
    return _rename(org, "initiatives", entity_id, body.name, "id, name, status")


@router.delete("/initiatives/{entity_id}")
def delete_initiative(entity_id: str, org: str = Depends(current_org("/initiatives"))) -> dict:
    # projects/applications reference initiative_id without ON DELETE — null them first.
    ctx, conn = _tx(org)
    try:
        conn.execute(text("update projects set initiative_id = null where initiative_id = :id"), {"id": entity_id})
        conn.execute(text("update applications set initiative_id = null where initiative_id = :id"), {"id": entity_id})
        gone = conn.execute(text("delete from initiatives where id = :id returning id"), {"id": entity_id}).mappings().all()
    finally:
        ctx.__exit__(None, None, None)
    if not gone:
        raise HTTPException(status_code=404, detail="initiative not found")
    return {"deleted": entity_id}


@router.patch("/projects/{entity_id}")
def rename_project(entity_id: str, body: RenameBody, org: str = Depends(current_org("/projects"))) -> dict:
    return _rename(org, "projects", entity_id, body.name, "id, name, slug, squad_id, repository_id")


@router.delete("/projects/{entity_id}")
def delete_project(entity_id: str, org: str = Depends(current_org("/projects"))) -> dict:
    # Cascades tasks/memory/gates; also drop this project's conversation transcripts.
    ctx, conn = _tx(org)
    try:
        conn.execute(text("delete from thread_transcript where project_id = cast(:id as uuid)"), {"id": entity_id})
        gone = conn.execute(text("delete from projects where id = :id returning id"), {"id": entity_id}).mappings().all()
    finally:
        ctx.__exit__(None, None, None)
    if not gone:
        raise HTTPException(status_code=404, detail="project not found")
    return {"deleted": entity_id}
